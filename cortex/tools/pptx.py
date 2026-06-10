"""
cortex.tools.pptx
──────────────────
Create PowerPoint (.pptx) presentations programmatically with python-pptx.

Requires python-pptx for the real file:
    pip install python-pptx

The model passes a list of slide specs; this tool handles layout, themes,
inline markdown (**bold**, *italic*, `code`), speaker notes and footer.

Promoted to a native cortex tool from a script cortex generated itself.
The rendering engine (create_presentation) is unchanged and battle-tested;
SCHEMA + execute() are the thin LLM-facing wrapper.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

SCHEMA = {
    "type": "function",
    "function": {
        "name": "pptx",
        "description": (
            "Create a formatted PowerPoint (.pptx) presentation. "
            "Use this for ANY request to 'make a presentation', 'create slides', "
            "'a deck', 'una presentación', or a .pptx file. "
            "Pass the output path and a list of slides. Each slide has a title and "
            "optional subtitle, content (bullet list), layout, and speaker notes. "
            "Bullet text supports inline markdown: **bold**, *italic*, `code`."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Output .pptx path. Example: ~/Desktop/deck.pptx",
                },
                "slides": {
                    "type": "array",
                    "description": "List of slide objects, in order.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide title (required)."},
                            "subtitle": {"type": "string", "description": "Optional subtitle line."},
                            "content": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Bullet points. Supports **bold**, *italic*, `code`.",
                            },
                            "layout": {
                                "type": "string",
                                "enum": ["title", "title_content", "section"],
                                "description": "Slide layout. 'section' = full-color divider slide.",
                            },
                            "notes": {"type": "string", "description": "Speaker notes (optional)."},
                        },
                        "required": ["title"],
                    },
                },
                "theme": {
                    "type": "string",
                    "enum": ["light", "dark", "corporate", "sunset"],
                    "description": "Visual theme. Default 'light'.",
                },
                "title": {
                    "type": "string",
                    "description": "Document title (metadata / author). Optional.",
                },
            },
            "required": ["path", "slides"],
        },
    },
}


# ── Themes ──────────────────────────────────────────────────────────────────────────

THEMES: dict[str, dict[str, str]] = {
    "light": {
        "bg": "FFFFFF", "title": "1F2937", "subtitle": "4B5563",
        "body": "111827", "accent": "3B82F6", "muted": "6B7280",
    },
    "dark": {
        "bg": "0F172A", "title": "F8FAFC", "subtitle": "CBD5E1",
        "body": "E2E8F0", "accent": "38BDF8", "muted": "94A3B8",
    },
    "corporate": {
        "bg": "F8FAFC", "title": "0B2545", "subtitle": "13315C",
        "body": "1F2937", "accent": "D4AF37", "muted": "64748B",
    },
    "sunset": {
        "bg": "FFF7ED", "title": "7C2D12", "subtitle": "9A3412",
        "body": "1C1917", "accent": "F97316", "muted": "78716C",
    },
}


@dataclass
class SlideSpec:
    title: str
    subtitle: str | None = None
    content: list[str] | None = None
    layout: str = "title_content"  # title, title_content, section
    image: str | None = None
    notes: str | None = None


# ── inline markdown ───────────────────────────────────────────────────────────────────

def _parse_inline_runs(text: str) -> list[tuple[str, dict]]:
    """Convert simple inline markdown to a list of (text, format) tuples."""
    runs: list[tuple[str, dict]] = []
    pos = 0
    pattern = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`")
    for m in pattern.finditer(text):
        if m.start() > pos:
            runs.append((text[pos:m.start()], {}))
        if m.group(1) is not None:
            runs.append((m.group(1), {"bold": True}))
        elif m.group(2) is not None:
            runs.append((m.group(2), {"italic": True}))
        else:
            runs.append((m.group(3), {"code": True}))
        pos = m.end()
    if pos < len(text):
        runs.append((text[pos:], {}))
    return runs


# ── slide painting ────────────────────────────────────────────────────────────────────

def _set_slide_background(slide, hex_color, RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _add_accent_bar(slide, theme, slide_w, MSO_SHAPE, RGBColor, Inches) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.18))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(theme["accent"])
    bar.shadow.inherit = False


def _add_title(slide, text, theme, slide_w, PP_ALIGN, RGBColor, Inches, Pt) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), slide_w - Inches(1.4), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string(theme["title"])


def _add_subtitle(slide, text, theme, slide_w, PP_ALIGN, RGBColor, Inches, Pt) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), slide_w - Inches(1.4), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string(theme["subtitle"])


def _add_bullets(slide, items: Iterable[str], theme, slide_w, PP_ALIGN, RGBColor, Inches, Pt) -> None:
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), slide_w - Inches(1.8), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0

    first = True
    for raw in items:
        text = re.sub(r"^\s*[-*]\s+", "", raw).strip()
        if not text:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.level = 0

        bullet = p.add_run()
        bullet.text = "•  "
        bullet.font.size = Pt(20)
        bullet.font.bold = True
        bullet.font.color.rgb = RGBColor.from_string(theme["accent"])

        for txt, fmt in _parse_inline_runs(text):
            if not txt:
                continue
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(20)
            r.font.name = "Calibri"
            r.font.color.rgb = RGBColor.from_string(theme["body"])
            if fmt.get("bold"):
                r.font.bold = True
            if fmt.get("italic"):
                r.font.italic = True
            if fmt.get("code"):
                r.font.name = "Consolas"
                r.font.size = Pt(18)


def _add_footer(slide, theme, slide_w, page, total, PP_ALIGN, RGBColor, Inches, Pt) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), slide_w - Inches(1.4), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Cortex  ·  {page} / {total}"
    run.font.size = Pt(10)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string(theme["muted"])


def _add_section_slide(slide, title, theme, slide_w, slide_h, MSO_SHAPE, PP_ALIGN, RGBColor, Inches, Pt) -> None:
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor.from_string(theme["accent"])

    box = slide.shapes.add_textbox(Inches(0), slide_h / 2 - Inches(0.7), slide_w, Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")


# ── public render API ─────────────────────────────────────────────────────────────────

def _normalize(slides: list) -> list[SlideSpec]:
    out: list[SlideSpec] = []
    for s in slides:
        if isinstance(s, SlideSpec):
            out.append(s)
        elif isinstance(s, dict):
            out.append(SlideSpec(
                title=s.get("title", ""),
                subtitle=s.get("subtitle"),
                content=s.get("content"),
                layout=s.get("layout", "title_content"),
                image=s.get("image"),
                notes=s.get("notes"),
            ))
        else:
            raise TypeError(f"Invalid slide: {s!r}")
    return out


def create_presentation(slides, output_path, theme: str = "light", title: str | None = None) -> Path:
    """Create a .pptx from a list of slide dicts/SlideSpecs. Returns the saved path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    if theme not in THEMES:
        raise ValueError(f"Theme '{theme}' does not exist. Options: {list(THEMES)}")

    pal = THEMES[theme]
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    slide_w, slide_h = prs.slide_width, prs.slide_height

    blank = prs.slide_layouts[6]
    specs = _normalize(slides)
    total = len(specs)
    if total == 0:
        raise ValueError("Slide list is empty.")

    if title:
        prs.core_properties.title = title
        prs.core_properties.author = "Cortex"

    for idx, spec in enumerate(specs, start=1):
        slide = prs.slides.add_slide(blank)
        _set_slide_background(slide, pal["bg"], RGBColor)

        if spec.layout == "section":
            _add_section_slide(slide, spec.title, pal, slide_w, slide_h,
                               MSO_SHAPE, PP_ALIGN, RGBColor, Inches, Pt)
        else:
            _add_accent_bar(slide, pal, slide_w, MSO_SHAPE, RGBColor, Inches)
            _add_title(slide, spec.title, pal, slide_w, PP_ALIGN, RGBColor, Inches, Pt)
            if spec.subtitle:
                _add_subtitle(slide, spec.subtitle, pal, slide_w, PP_ALIGN, RGBColor, Inches, Pt)
            if spec.content:
                _add_bullets(slide, spec.content, pal, slide_w, PP_ALIGN, RGBColor, Inches, Pt)
            if spec.image and Path(spec.image).exists():
                slide.shapes.add_picture(spec.image, Inches(8.5), Inches(4.5), height=Inches(2.3))

        if spec.notes:
            slide.notes_slide.notes_text_frame.text = spec.notes

        if spec.layout != "section":
            _add_footer(slide, pal, slide_w, idx, total, PP_ALIGN, RGBColor, Inches, Pt)

    prs.save(str(output_path))
    return output_path


# ── LLM-facing executor ───────────────────────────────────────────────────────────────

def _coerce_slides(slides) -> list:
    """Models may pass slides as a real list or as a JSON string. Accept both."""
    if isinstance(slides, str):
        slides = json.loads(slides)
    if isinstance(slides, dict):  # {"slides": [...]} or a single slide
        slides = slides.get("slides", [slides])
    if not isinstance(slides, list):
        raise TypeError("'slides' must be a list of slide objects.")
    return slides


def execute(path: str, slides, theme: str = "light", title: str = "", **_) -> str:
    try:
        from pptx import Presentation  # noqa: F401 — fail fast if missing
    except ImportError:
        return (
            "[ERROR] python-pptx not installed.\n"
            "Install: pip install python-pptx"
        )
    try:
        specs = _coerce_slides(slides)
        out = create_presentation(specs, path, theme=theme or "light", title=title or None)
        kb = out.stat().st_size // 1024 + 1
        return f"Created {out}  ({len(specs)} slides, theme '{theme or 'light'}', {kb} KB)"
    except Exception as e:
        return f"[ERROR] pptx: {type(e).__name__}: {e}"
